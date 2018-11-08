# -*- coding: utf-8 -*-
"""
some basic functions used to generate PowerPoint report
@author: junxu
"""

from numpy import unique, nonzero, shape, mean, std, zeros,ones, arange, mod, asarray,c_,squeeze, round,dot

from scipy.stats.mstats import normaltest
from scipy.stats import f_oneway, shapiro, kstest,anderson,ks_2samp, ttest_ind
from scipy.stats.mstats import kruskalwallis
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from datetime import date
from os import listdir
from os.path import isfile, join
import matplotlib.pyplot as plt
from datetime import date
from DateTime import DateTime
import glob
import shutil
import os

def save_file_list(dirS,dirD=None, extName=None):
    '''
    copy file with specified extension from dirS to dirD
    dirS: source file directory
    dirD: destination directory
    extName: file extension
    '''
    
    if dirS[-1]!="/" or dirS[-2:]!="\\":
        dirS=dirS+"/"
        
    if dirD[-1]!="/" or dirD[-2:]!="\\":
        dirD=dirD+"/"
        
    files=(glob.glob(dirS+'*.'+extName))
    #print(files)
    
        
    for filename in files:
        
        filebase = os.path.basename(filename)
    #    ext = os.path.splitext(filename)[1]
    #    name=filebase.replace(ext,'')    
        shutil.copy2(filename, (dirD+filebase))

class options_class:
  def __init__(self,save_figure=1,plot_figure=1, plot_fontsize=10, lba_size_idx=50, access_type=2,report_name='trace_analysis_raw.ppt',time_interval=50,S2_threshold=1024,S2_threshold2=64,lba_size_set=50,max_stream_length=1024,seq_size_threshold=1024,ppt_template='Workload.pptx'):
     # self.lba_size_dist = lba_size_dist
     self.lba_size_idx = lba_size_idx
     self.save_figure=save_figure
     self.plot_figure=plot_figure
     self.plot_fontsize=plot_fontsize
     self.queue_len_setting=2.0 ** (arange(0,8,1))
     self.access_type=access_type
     self.report_name=report_name
     self.time_interval=time_interval
     self.S2_threshold=S2_threshold
     self.S2_threshold2=S2_threshold2
     self.lba_size_set=lba_size_set
     self.max_stream_length=max_stream_length
     self.seq_size_threshold=seq_size_threshold
     self.section_name=''
     self.ppt_template="Workload.pptx"
     self.ppt_name='analysis_report.pptx'
     self.lba_size_set=50


#    test_descriptions = {
#        'rand_4'  : (1, 'Replica 1 Read'),
#        'rand_2'  : (2, 'Replica 3 Read'),
#        'rand_6'  : (3, 'EC (6,3) Read'),
#        'write_3' : (4, 'Replica 1 Write'),
#        'write_1' : (5, 'Replica 3 Write'),
#        'write_5' : (6, 'EC (6,3) Write')
#    }


class slide_info_class:
    def __init__(self, width=4.5, height=3,top=1.3,left=1,font_size=20):
        self.width = width
        self.height = height
        self.top=top
        self.left=left
        self.font_size=font_size
        
def generate_1fig_page(prs,filenames,str_list, content, locations=slide_info_class()):
    '''
    generate a slide with 2 figures in 3 axes + 1 full length text
    Input:
        filenames: a list of picture filenames
        str_list: list of strings for title and subtitle
        content: the string list for comments
    '''    
    width=locations.width
    height=locations.height
    width0=Inches(width)
    height0=Inches(height)
    top=locations.top
    left=locations.left
    tops=[top, top, top+height]
    lefts=[left, left+width, left+width]
    font_size=locations.font_size
    
    out_slide_layout = prs.slide_layouts[12]
    out_slide = prs.slides.add_slide(out_slide_layout)
    title = out_slide.shapes.title
    title.text = str_list[0]
    subtitle = out_slide.placeholders[13]    
    tf = subtitle.text_frame 
    tf.text=str_list[1]
    pic_name=filenames[0]
    pic = out_slide.shapes.add_picture(pic_name, Inches(left), Inches(top), height=height0, width=width0)     
    subtitle = out_slide.placeholders[12] 
    subtitle.text=''  
    add_items_to_tfl2(out_slide, content,font_size,place_no=12)           
        
def generate_2fig_page(prs,filenames,str_list, content, locations=slide_info_class()):
    '''
    generate a slide with 2 figures in 3 axes + 1 full length text
    Input:
        filenames: a list of picture filenames
        str_list: list of strings for title and subtitle
        content: the string list for comments
    '''    
    width=locations.width
    height=locations.height
    width0=Inches(width)
    height0=Inches(height)
    top=locations.top
    left=locations.left
    tops=[top, top, top+height]
    lefts=[left, left+width, left+width]
    font_size=locations.font_size

    out_slide_layout = prs.slide_layouts[11]
    out_slide = prs.slides.add_slide(out_slide_layout)
    title = out_slide.shapes.title
    title.text = str_list[0]
    subtitle = out_slide.placeholders[15]    
    tf = subtitle.text_frame 
    tf.text=str_list[1]
    for j in arange(2).reshape(-1):
        pic_name=filenames[j]
        pic = out_slide.shapes.add_picture(pic_name, Inches(lefts[j]), Inches(tops[j]), height=height0,width=width0)         
    subtitle = out_slide.placeholders[13] 
    subtitle.text=''
    
    add_items_to_tfl2(out_slide, content,font_size,place_no=13)
        

def generate_3fig_page(prs,filenames,str_list, content, locations=slide_info_class()):
    '''
    generate a slide with 3 figures in 3 axes + 1 half length text
    Input:
        filenames: a list of picture filenames
        str_list: list of strings for title and subtitle
        content: the string list for comments
    '''    
    width=locations.width
    height=locations.height
    width0=Inches(width)
    height0=Inches(height)
    top=locations.top
    left=locations.left
    tops=[top, top, top+height]
    lefts=[left, left+width, left+width]
    font_size=locations.font_size

    out_slide_layout = prs.slide_layouts[10]
    out_slide = prs.slides.add_slide(out_slide_layout)
    title = out_slide.shapes.title
    title.text = str_list[0]
    subtitle = out_slide.placeholders[15]    
    tf = subtitle.text_frame 
    tf.text=str_list[1]
    for j in arange(3).reshape(-1):
        pic_name=filenames[j]
        pic = out_slide.shapes.add_picture(pic_name, Inches(lefts[j]), Inches(tops[j]), height=height0,width=width0) 
    add_items_to_tf2(out_slide, content,font_size,place_no=14)
  
    
def generate_4fig_page(prs,filenames,str_list, content,locations=slide_info_class()):
    '''
    generate a slide with 3 figures in 3 axes
    Input:
        filenames: a list of picture filenames
        str_list: list of strings for title and subtitle
        content: the string list for comments
    '''
    width=locations.width
    height=locations.height
    width0=Inches(width)
    height0=Inches(height)
    top=locations.top
    left=locations.left
    tops=[top, top, top+height, top+height]
    lefts=[left, left+width, left, left+width]
    font_size=locations.font_size

    out_slide_layout = prs.slide_layouts[10]
    out_slide = prs.slides.add_slide(out_slide_layout)
    title = out_slide.shapes.title
    title.text = str_list[0]
    subtitle = out_slide.placeholders[15]    
    tf = subtitle.text_frame 
    tf.text=str_list[1]
    for j in arange(4).reshape(-1):
        pic_name=filenames[j]
        pic = out_slide.shapes.add_picture(pic_name, Inches(lefts[j]), Inches(tops[j]), height=height0,width=width0) 
   

def check_sequence(seek_value,seek_queue,options=None,*args,**kwargs):
    '''
 function str0=check_sequence(seek_value, seek_queue, options)
 check the sequence/randomness based on pre-defined threshold.
 options.seek [4]: if <[1], low; if 1~[2], relative low; if 2~[3] relatively
 high; if >[3] high
 seek_delta: 1: the threshold for steady state; 2/3: the treshold for mixed stream 
 options.seek_head_str: string to add at the head
 seek_value: nXm matrix; m is usually equal to 10::: # 1 total R/W IO
 number, 2 sequnce number, 3 mean, 4 mean abs, 5 median, 6 mode, 7 mode
 couter, 8 min abs, 9 max abs, 10 std abs  
 seek_queue: nX1 vector, queue length
     '''    
    str0=[]
    a,b=shape(seek_value)
    a1=len(seek_queue)
    # print([a,a1])
    if a != a1:
        disp('Error in input values dimension')
        return str0
    
    if hasattr(options,'seek_head_str'):
        seek_head_str=options.seek_head_str
    else:
        seek_head_str=''
    
    if hasattr(options,'seek'):
        seek=options.seek
    else:
        seek=[0.1,0.2,0.5,0.8]
    
    if hasattr(options,'seek_delta'):
        seek_delta=options.seek_delta
    else:
        seek_delta=[0.01,0.1,0.2]
    
    if a < 2:
        disp('Warning! too small size of value! extend the quene length and run again')
        return str0
    
    # only when mode==0, we consider the trace has seqnence
    temp_str0=(('QL='+str(seek_queue[0])))
    if seek_value[0,5] == 0:
        mode_rati=seek_value[0,1] / seek_value[0,0]
        if mode_rati < seek[0]:
            temp_str1='Very low'
        else:
            if mode_rati < seek[1]:
                temp_str1='Low'
            else:
                if mode_rati < seek[2]:
                    temp_str1='Relatively high'
                else:
                    if mode_rati < seek[3]:
                        temp_str1='High'
                    else:
                        if mode_rati < seek[3]:
                            temp_str1='Very high'
                        else:
                            temp_str1='Unknown'
        temp_str1=(('Mode=0 with ratio= '+"{0:.3f}".format(mode_rati)+' and sequence '+temp_str1+'  at '+temp_str0))
        # first check the steady state
        if a < 3:
            disp('Warning! too small size of value to check steady state! steady status will not be included')
            temp_str2=[]
        else:
            steady_i=0
            for i in arange(a-1).reshape(-1):
                if seek_value[i + 1,3] / seek_value[i,3] < dot(seek_delta[0],(seek_queue[i + 1] / seek_queue[i])):
                    steady_i=seek_queue[i + 1]
                    break
            if steady_i > 0:
                temp_str2=(('Steady state possibly at QL='+str((steady_i))))
            else:
                temp_str2=(('Higher rate possibly for long length than '+str(seek_queue[a-1])))
            for i in arange(a - 1).reshape(-1):
                if seek_value[i + 1,3] / seek_value[i,3] > dot(seek_delta[2],(seek_queue[i + 1] / seek_queue[i])):
                    temp_str3='Strong mixed streams detected'
                    break
                else:
                    if seek_value[i + 1] / seek_value[i] > dot(seek_delta[1],(seek_queue[i + 1] / seek_queue[i])):
                        temp_str3='Mixed streams detected'
                    else:
                        temp_str3=''
    else:
        temp_str1=(('Mode='+str(seek_value[0,5])+'  at '+temp_str0+'; Relatively random'))
        temp_str3=''
        temp_str2=''
    
    str0=((temp_str1+'; '+temp_str2+'; '+temp_str3))
    
    return str0


def parse_time(time_list):
    
    time_array=zeros((len(time_list),1))
    j=-1;
    for i in time_list:
        j=j+1        
        dt=DateTime(i)
        time_array[j]=dt.timeTime()
    return time_array    

  

def generate_table_page(prs, V, str_list,x_list,y_list,options=1,create_page=1,shapes=None, location=[0.6,1.3]):
    '''
    prs: object of the presentation
    V: matrix to create the table; the dimension must be consistent to x_list and y_list, i.e., len(x_list)==shape(V)[0] && len(y_list)==shape(V)[1]
    str_list: string list of some title/item names
    x_list: name list for x-axis except for (0,0)
    y_list: name list for y-axis except for (0,0)
    options: 0: only print out the value of V; 1: a simple analysis for the value V with mean and std
    '''
    
    if len(x_list)!=shape(V)[0] or len(y_list)!=shape(V)[1]:
        print('Error! the dimension of the data array does not match x/y list')
        return -1
    
    if create_page==1:
        cpu_slide_layout = prs.slide_layouts[15]
        # Create the summary graph
        slide = prs.slides.add_slide(cpu_slide_layout)
        title = slide.shapes.title
        title.text = str_list[0]    #"Network Performance"
        placeholder = slide.placeholders[11]
        shapes = slide.shapes
        subtitle = slide.placeholders[14]
        subtitle.text = str_list[1]
    
    
    if options==0:
        cols = len(y_list)+1
        rows = len(x_list)+1
    else:
        cols = len(y_list)+1
        rows = len(x_list)+3
    
    left = Inches(location[0])
    top = Inches(location[1])
    width = int(prs.slide_width*0.9) #Inches(10)
    width_s=int((width-left)/cols)
    height = int(prs.slide_height-Inches(1.5)) #Inches(7)
    height_s=int((height-top-Inches(1.5))/rows)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    # set column widths
    if cols>20:
        w0=0.5
    elif cols>10:
        w0=1
    elif cols>5:
        w0=1.5
    else:
        w0=1.8
            
    table.columns[0].width = Inches(w0)
    table.rows[0].height = Inches(0.5)
#    table.rows[1].height = Inches(1.0)
#    table.rows[2].height = Inches(1.0)
    table.cell(0, 0).text=str_list[2] #'Network'
       
    for i in arange(0,cols-1):
        table.columns[i+1].width = (width_s)
        
    for i in arange(0,rows-1):
        table.rows[i+1].height = (height_s)

    for i in arange(0,cols-1):
        table.cell(0, i+1).text = y_list[i]
        
    if options>0:
        rows0=rows-2
        table.cell(rows-2, 0).text='mean'
        table.cell(rows-1, 0).text='std'
        V_s=zeros((2,cols-1))
        for i in arange(0,cols-1):
           V_s[0,i]=mean(V[:,i]) 
           V_s[1,i]=std(V[:,i])
        for j in arange(0,2): # (idx-1)*2,(idx-1)*2+1
            for i in arange(1,cols):
                if V_s[j,i-1]==0:
                    table.cell(j+rows0, i).text = "{0:.0f}".format(V_s[j,i-1])
                elif (V_s[j,i-1]==1) or (V_s[j,i-1]==--1):
                    table.cell(j+rows0, i).text = "{0:.0f}".format(V_s[j,i-1])
                else:
                    table.cell(j+rows0, i).text = "{0:.3f}".format(V_s[j,i-1]) #str(V[j-1,i-1])
    else:
        rows0=rows
        
        
    for i in arange(rows0-1):
        table.cell(i+1, 0).text=str(x_list[i]) #'Speed'
        
        
    for j in arange(1,rows0): # (idx-1)*2,(idx-1)*2+1
        for i in arange(1,cols):
            print([j,i])
            if V[j-1,i-1]==0:                
                table.cell(j, i).text = "{0:.0f}".format(V[j-1,i-1])
            #elif (V[j-1,i-1])==-1 or (V[j-1,i-1])==1 or (V[j-1,i-1])==-2:
            elif (V[j-1,i-1]).is_integer():
                table.cell(j, i).text = "{0:.0f}".format(V[j-1,i-1])                
            else:
                table.cell(j, i).text = "{0:.3f}".format(V[j-1,i-1]) #str(V[j-1,i-1])
       
    font_size=item_font_size(max([cols,rows]))
             
    for j in arange(0,rows): # (idx-1)*2,(idx-1)*2+1
        for i in arange(0,cols):   
            cell = table.rows[j].cells[i]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(font_size)
    if create_page==1:
        return slide
    else:
        return -1

def item_font_size(sub_lines):
    
    if sub_lines>30:
        font_size=5
    elif sub_lines>25:
        font_size=6          
    elif sub_lines>20:
        font_size=8        
    elif sub_lines>12:
        font_size=10
    elif sub_lines>10:
        font_size=11    
    elif sub_lines>8:
        font_size=12
    elif sub_lines>6:
        font_size=14
    else:
        font_size=16
    return font_size    

def add_items_to_tf(out_slide, text_str, test_list,font_size,total_num=-1,place_no=10):
    '''
    add the listed items in test_list into a text frame of a placeholder in two levels
    '''
    subtitle = out_slide.placeholders[place_no]
    test_num=len(test_list)
    tf = subtitle.text_frame
    if tf.text!='':
        p = tf.add_paragraph()
        p.text = text_str # +":"+str(test_num)
    else:
        tf.text = text_str #+":"+str(test_num)
    for i0 in arange(test_num):
        i=test_list[i0]
        p = tf.add_paragraph()
        p.text = i
        p.level=1 
        p.font.size=Pt(font_size)

def add_items_to_tf2(out_slide, test_list,font_size,place_no=10):
    '''
    add the listed items in test_list into a text frame of a placeholder in one level
    '''
    subtitle = out_slide.placeholders[place_no]
    test_num=len(test_list)
    tf = subtitle.text_frame
    if tf.text!='':
        p = tf.add_paragraph()
        p.text = test_list[0] # +":"+str(test_num)
    else:
        tf.text = test_list[0] #+":"+str(test_num)
    for i0 in arange(test_num-1):
        i=test_list[i0+1]
        p = tf.add_paragraph()
        p.text = i
#        p.font.size=Pt(font_size)

def add_items_to_tfl(out_slide,text_str, test_list,font_size,total_num=-1,place_no=10):
    '''
    add the listed items in test_list into a text frame of a placeholder in two levels with multiple items
    text_str, test_list must be the list of strings; not a string; use [] if empty
    the length of text_str shall be equal to that of test_list
    the item in test_list is lower level than these item in text_str
    '''
    subtitle = out_slide.placeholders[place_no]
    test_num=len(test_list)
    text_num=len(text_num)
    if test_num!=text_num:
        print('Error! the format the string list is wrong. Please check the dimention')
    tf = subtitle.text_frame

    for i in arange(test_num):    
        if tf.text!='':
            p = tf.add_paragraph()
            p.text = text_str[i] # +":"+str(test_num)
        else:
            tf.text = text_str[i] #+":"+str(test_num)
        for i0 in arange(len(test_num[i])):
            i=test_list[i][i0]
            p = tf.add_paragraph()
            p.text = i
            p.level=1 
            #p.font.size=Pt(font_size)

def add_items_to_tfl2(out_slide, test_list,font_size,total_num=-1,place_no=10):
    '''
    add the listed items in test_list into a text frame of a placeholder in two levels with multiple items
    input string list format is different from add_items_to_tfl
    text_str, test_list are combined into test_list, which must be even length. One at level 0 followed by a list of  level 1 items, e.g.,
    test_list=['level 0', ['level 1', 'level 1'], 'level 0', [], 'level 1', ['level 1']]
    '''
    subtitle = out_slide.placeholders[place_no]
    test_num=len(test_list)
    if mod(test_num,2)!=0:
        print('Error! the format the string list is wrong. Please check the length, which must be even')
    tf = subtitle.text_frame

    for i in arange(int(test_num/2)):    
        text=test_list[(i)*2]
        text2=test_list[(i)*2+1]
        if tf.text!='':
            p = tf.add_paragraph()
            p.text = text # +":"+str(test_num)
        else:
            tf.text = text #+":"+str(test_num)
        for i0 in arange(len(text2)):
            i=text2[i0]
            p = tf.add_paragraph()
            p.text = i
            p.level=1 
            #p.font.size=Pt(font_size)

def two_items_per_line(test_list):
    '''
    form 2 items in the test_list as one item; if odd, leave the last one as one item
    '''
    test_num=len(test_list)
    div_num=test_num/2
    new_list=[]
    for i in arange(div_num):
        new_list.append(test_list[i*2]+';'+test_list[i*2+1])
    if mod(test_num,2)==1:
        new_list.append(test_list[test_num-1])     
    return new_list        
    
def all_items_in_one_line(test_list,options=0):
    #test_num=len(test_list)
    str0=''
    j=-1
    for i in test_list:
        j=j+1;
        if options==1:
            str0=str0+'('+str(j)+') '+i+', '
        else:
            str0=str0+i+', '
    return str0[0:-2]       


def get_rados_bw(round_list):
    test_list,valid_data,sum_data,raw_data=import_all_jsons(round_list)
    Con_array, Con_array2, Con_array3, consistent_list, consistent_list2, consistent_list3, consistent_list4   =statistical_comparison(round_list, test_list,raw_data,valid_data,sum_data)
    return Con_array,Con_array2, consistent_list


def check_value_consistency(v_list, ck_pt=10, th=0.2):
    '''
    check if the values in the given lists are similar to each others at given time windows
    v_list is a nxD list
    ck_pt is the number of check points
    th is the thresholds to indicate the tolerance range, th[0] for mean, and th[1] for std/mean
    '''
    if shape(v_list)[1]<ck_pt:
        print('the number of check points is larger than the total number of data points; choose the smaller number ck_pt')
        ck_pt=shape(v_list)[1]
    j1=0;
    interval=shape(v_list)[1]/1.0/ck_pt
    ck_list=zeros((ck_pt,1))
    l0=len(v_list)
    for i in arange(ck_pt):
        j2=int(floor((i+1)*interval))
        m0=zeros((l0,2))
        for j in arange(l0):
            vt=v_list[j][j1:j2]
            m0[j,0]=mean(vt)
            m0[j,1]=std(vt)
        m1=mean(m0[:,0])
        m2=std(m0[:,0])
        if abs(m2/m1)<=th:
            ck_list[i]=1
    con_ratio=sum(ck_list)/ck_pt    
    return con_ratio

def check_value_consistency2(v_list, th=0.2):
    '''
    check if the values in the given lists are similar to each others at given time windows
    v_list is a nxD list
    ck_pt is the number of check points
    th is the thresholds to indicate the tolerance range, th[0] for mean, and th[1] for std/mean
    '''
    l0=len(v_list)
    v_mean=mean(v_list)
    v_up=v_mean*(1+th)
    v_down=v_mean*(1-th)
    con0=0
    for i in arange(l0):
        if (v_list[i]>=v_down) and (v_list[i]<=v_up):
            con0=con0+1
    con_ratio=(1.0*con0)/l0            
    return con_ratio

#str_list=['Rados Bandwidth', 'Individual rounds', 'Thoughput (MBPS)']

def generate_figure_page(prs, V_df,str_list,key_list,time_array=[]):
    
    slide_layout = prs.slide_layouts[12]
    # Create the summary graph
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = str_list[0]    # slide title
    subtitle = slide.placeholders[13] # sub title
    subtitle.text = str_list[1]   
    x=len(V_df)
    fig=plt.figure()
    ax=fig.add_subplot(1,1,1)
    
    for i in arange(x):   
        if len(time_array)==0:         
            ax.plot(arange(len(V_df[i])),V_df[i],label=str(key_list[i])) 
        else:
            ax.plot(time_array,V_df[i],label=str(key_list[i])) 

    handles, labels = ax.get_legend_handles_labels()
    ax.legend(handles, labels,loc=0,fontsize='small')  
    ax.set_ylabel(str_list[2])
    ax.set_xlabel('Time')
    ax.set_title(str_list[0])
    pic_name='temp.png'
    plt.savefig(pic_name)                                                                                             
    
    # add to the figure to slide 
    im=Image.open(pic_name,'r')
    iw,ih=im.size # (width,height) tuple
    top0=Inches(1.1)
    left0=Inches(2)
    pic = slide.shapes.add_picture(pic_name, left0, top0)     
    
    str_array2=[]
    
    str0='The mean value:'
    for i in arange(x):
        str0=str0+(key_list[i])+'('+str(round(mean(V_df[i]),2))+'); '
    str_array2.append(str0)    
    
    str0='The ratio within +- 20% of mean:'
    for i in arange(x):
        rate=check_value_consistency2(V_df[i],0.2)
        str0=str0+(key_list[i])+'('+str(round(rate,2))+'); '
    str_array2.append(str0)
    
    str0='The ratio within +- 10% of mean:'
    for i in arange(x):
        rate=check_value_consistency2(V_df[i],0.1)
        str0=str0+(key_list[i])+'('+str(round(rate,2))+'); '
    str_array2.append(str0)    
    
    subtitle = slide.placeholders[12] # sub title
    subtitle.text = ''       
    add_items_to_tf2(slide, str_array2,font_size=10, place_no=12)

def add_metric_to_slide(prs,round_list,path0,metric_str,options=[1,1,1,1,1,1],gen_sum=0,reduced_num=0,method_v=0):      
    '''
    generate slides to presentation with given options
    Input parameters:
        prs: the handle of presenation object
        round_list: the json filename list of all test rounds
        path0: the path contains the files
        metric_str: string lists for description of the slides
        options: a vector specify if the analytical items (normal test for all rounds, average value for all rounds, pair comparison pages based on normal distribution,pair comparison pages based on continue distribution,  pair comparison pages for mean) should be output or not --> 0/no ouput; 1/output. 
    '''
    #print(path0)    
    #print(round_list[0])
    test_list,valid_data,sum_data,raw_data,time_data=import_all_jsons(round_list,path0,metric_str[0],metric_str[1],gen_sum,reduced_num,method_v)
    Con_array, Con_array2, Con_array3, consistent_list, consistent_list2, consistent_list3, consistent_list4  =statistical_comparison(round_list, test_list,raw_data,valid_data,sum_data)
    
    test_num=len(test_list)
    round_num=len(round_list)
    
    data_set_val=zeros((test_num,3))
    
    for tj in arange(test_num):
        idx1=shape(nonzero(valid_data[:,tj,0]==1))[1]
        idx0=shape(nonzero(valid_data[:,tj,0]==0))[1]
        idx2=round_num-idx1-idx0
        data_set_val[tj,:]=[idx1,idx0,idx2]
        
    
    len_cl=len(consistent_list)
    ck_cl=zeros((len_cl,1))
    ck_cl2=zeros((len_cl,1))
    for i in arange(len_cl):
        if (consistent_list[i][1]<0) :
            ck_cl[i]=-1
        elif consistent_list[i][1]>0.05:
            ck_cl[i]=1
        if (consistent_list3[i][1]<0) :
            ck_cl2[i]=-1
        elif consistent_list3[i][1]>0.05:
            ck_cl2[i]=1    
            
            
    # overall summary page
    str_list=[metric_str[1],'Benchmark Consistency Summary for all Rounds','value']       
    x_list=['normal','non-normal','invalid','F-value','p-value','F-result','H-value','p-value','H-result']
    slide=generate_table_page(prs, asarray(c_[data_set_val, squeeze(consistent_list)[:,0:2],ck_cl,squeeze(consistent_list3)[:,0:2],ck_cl2]).T, str_list,x_list,test_list,0)
    
    title = slide.placeholders[13]  
    title.text = '1: likely same mean/median; 0: unlikely same mean/median; -1: no enough data for test'
    
    round_list2=[]
    for i in arange(round_num):
        round_list2.append('R'+str(i))         
    
    if options[0]==1:
        # normal test for all rounds
        str_list=[metric_str[1],'Normal test Summary for all Rounds/tests','normal']
        slide=generate_table_page(prs, (valid_data[:,:,0]).T, str_list,test_list,round_list2,0)
        
        title = slide.placeholders[13]  
        title.text = '1: likely normal distribution; 0: unlikely normal distribution; -1: no enough data for test'
 
    if options[1]==1:
        # average value for all rounds
        str_list=[metric_str[1],'Average value for all Rounds/tests','mean']
        slide=generate_table_page(prs, (valid_data[:,:,1]), str_list,round_list2,test_list,1)
        
        title = slide.placeholders[13]  
        title.text = metric_str[2]    
        
        # the ratio of value difference to mean of all
        V0=asarray(valid_data[:,:,1])
        mean_value=mean(V0,0)
#        print('mean') 
#        print(mean_value)
#        print('V0')
#        print(V0)
        V1=[]
        for i in arange(test_num):
            Vt=(V0[:,i]-mean_value[i])/mean_value[i]
#            print(V0)
#            print(V0[:,i]) 
#            print(Vt)
            V1.append(Vt)
        V1=asarray(V1)
#        print(V1)
#        print(shape(V1))
        str_list=[metric_str[1],'Ratio of difference from overall mean for all Rounds/tests','ratio']
        slide=generate_table_page(prs, V1.T, str_list,round_list2,test_list,1)
        
        title = slide.placeholders[13]  
        title.text = metric_str[2] 
        
        # the ratio of the value within a certain range
        V_ratio=zeros((round_num,test_num))
        for i in arange(round_num):
            time_array=[]
            for j in arange(test_num):
                v0=raw_data[i,j]
                V_ratio[i,j]=check_value_consistency2(v0, th=0.1)
        str_list=[metric_str[1],'Frequency ratio of values within +-10% of mean for individual round/test','ratio']
        slide=generate_table_page(prs, V_ratio, str_list,round_list2,test_list,1)
        
        # the ratio of the value within a certain range
        V_ratio2=zeros((round_num,test_num))
        for i in arange(round_num):
            time_array=[]
            for j in arange(test_num):
                v0=raw_data[i,j]
                V_ratio2[i,j]=check_value_consistency2(v0, th=0.2)
        str_list=[metric_str[1],'Frequency ratio of values within +-20% of mean for individual round/test','ratio']
        slide=generate_table_page(prs, V_ratio2, str_list,round_list2,test_list,1)        
        
    # pair comparison pages based on normal distribution
    if options[2]==1:
        for j in arange(test_num):
            str_list=[metric_str[1]+' of '+test_list[j],'Benchmark Consistency for Pairs of Rounds (Normal)',test_list[j]] 
            slide=generate_table_page(prs, Con_array[j], str_list,round_list2,round_list2,0)
            title = slide.placeholders[13]  
            title.text = '1: likely same mean; 0: unlikely same mean; -1: at least one vector no enough data for test; -2: at least one vector unlikely normal distribution'

    # pair comparison pages based on continue distribution
    if options[3]==1:
        for j in arange(test_num):
            str_list=[metric_str[1]+' of '+test_list[j],'Benchmark Consistency for Pairs of Rounds (Continuous)',test_list[j]] 
            slide=generate_table_page(prs, Con_array2[j], str_list,round_list2,round_list2,0)
            title = slide.placeholders[13]  
            title.text = '1: likely same continuous distribution; 0: unlikely same continuous distribution; -1: at least one vector no enough data for test;'

    # pair comparison pages for mean
    if options[4]==1:
        for j in arange(test_num):
            str_list=[metric_str[1]+' of '+test_list[j],'Benchmark Consistency for Pairs of Rounds (Mean)',test_list[j]] 
            slide=generate_table_page(prs, Con_array3[j], str_list,round_list2,round_list2,0)
            title = slide.placeholders[13]  
            title.text = '1: likely same mean; 0: unlikely same mean; -1: at least one vector no enough data for test;'

    # plot individual runs' curve with summary
    if options[5]==1:
        str_list=[metric_str[1], 'Individual rounds:', metric_str[2]]
        for i in arange(round_num):
            V_df=[] 
            time_array=[]
            for j in arange(test_num):
                V_df.append(raw_data[i,j])                
                #t0=parse_time(time_data[i,j])
                #time_array.append(t0)
            str_list[1]='Individual rounds:'+ round_list[i]   
            generate_figure_page(prs, V_df,str_list,test_list,time_array=[])
    
    return raw_data            


def add_bullet_slide(prs,str_list,font_size=20):
    '''
    # Create the bullet items into a slide using a string list
    '''
    sum_slide_layout = prs.slide_layouts[4]

    slide = prs.slides.add_slide(sum_slide_layout)
    title = slide.shapes.title
    title.text = str_list[0]

    body_shape = slide.placeholders[10]
    tf = body_shape.text_frame
    tf.text = str_list[1]
    
    if len(str_list)>=3:
        for i in str_list[2:]:
            p = tf.add_paragraph()
            p.text = i
#            p.font.size = Pt(font_size)    
    



def generate_report(round_list=[], path0=[],options=[1,1,1,1,1,1]): 
    '''
    generate the comparison report for the given json files
    input parameters:
        round_list: the list of json files names; 
        path0: the path that contains the json files. if round_list gives the full path, path0 is not required
    '''
       
    if len(round_list)==0: 
        if len(path0)==0:
            print('ERROR! You must enter either the file list or the folder/path contains the files!')
        else:
            round_list = [f for f in listdir(path0) if isfile(join(path0, f))]
      
    ## start to create slides
    prs = Presentation('Workload.pptx')
    
    slide_width=9.0;
    slide_height=6.5;
    
    ## create a title page
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    title.text = 'RADOS Test Result Comparison'
    subtitle = slide.placeholders[13]
    subtitle.text = "{:%m-%d-%Y}".format(date.today())
    
    subtitle = slide.placeholders[10]
    subtitle.text = "Generated by PBJ Parser and Analyzer \nWDLABS"
    
    ## create an outline page for all settings
    out_slide_layout = prs.slide_layouts[4]
    # Create the summary graph
    out_slide = prs.slides.add_slide(out_slide_layout)
    title = out_slide.shapes.title
    title.text = "Outline"
    subtitle = out_slide.placeholders[11]
    subtitle.text = "PBJ Rados Benchmark Settings"  
        
    sub_lines=0
    
    pbjtr=pbj.PbjTestReport(path0+round_list[0])
    test_list=pbjtr.tests.keys()
    
    round_num=len(round_list)
    test_num=len(test_list)    
    
    if test_num>=5:
        sub_lines=2
        new_list=all_items_in_one_line(test_list)
        new_list=[new_list]
    elif test_num>=3:
        sub_lines=sub_lines+test_num/2
        new_list=two_items_per_line(test_list)        
    else:
        sub_lines=sub_lines+test_num
        new_list=test_list
     
    if round_num>20:
        sub_lines=10
        new_list2=all_items_in_one_line(round_list,options=1)
        new_list2=[new_list2]
    elif round_num>10:        
        sub_lines=sub_lines+round_num/2
        new_list2=two_items_per_line(round_list)  
    else:
        sub_lines=sub_lines+round_num      
        new_list2=round_list    
    
    font_size=item_font_size(sub_lines) 
    
    subtitle = out_slide.placeholders[10]    
    tf = subtitle.text_frame 
    tf.text=''    
       
    add_items_to_tf(out_slide, 'Benchmark test(s):'+str(test_num), new_list,font_size)        
    
    add_items_to_tf(out_slide, 'Benchmark round(s):'+str(round_num), new_list2,font_size)
    
    # create a notation page
    out_slide_layout = prs.slide_layouts[4]
    # Create the summary graph
    out_slide = prs.slides.add_slide(out_slide_layout)
    title = out_slide.shapes.title
    title.text = "Notations"
    subtitle = out_slide.placeholders[11]
    subtitle.text = "PBJ Rados Benchmark Result Summary"   
    
    subtitle = out_slide.placeholders[10]    
    tf = subtitle.text_frame 
    tf.text=''    
    
    font_size=16   
    str_list=['normal: the number of rounds that the values satisfy  normal distribution',
          'non-normal: the number of rounds that the values do not satisfy  normal distribution',
          'invalid: the collected dataset is too small for a valid hypothesis  test',
          'F-value /p-value: F value/p-value returned from 2-sided F-test',
          'F-result: if p-value is less than 0.05, the hypothesis of same mean of normal distribution is rejected (0); if it is greater, then the hypothesis is not rejected (1). Otherwise, no enough data for F-test (-1)',
          'H-value/p-value: H value/p-value returned from one-way H-test',
          'H-result: if p-value is less than 0.05, the hypothesis of same median is rejected (0); if it is greater, then the hypothesis is not rejected (1). Otherwise, no enough data for H-test (-1)']            
            
    
    add_items_to_tf(out_slide, 'Consistency summary for all rounds', str_list,font_size)
    
    str_list=['1: likely normal distribution',
              '0: unlikely normal distribution', 
              '-1: no enough data for test']            
    
    add_items_to_tf(out_slide, 'Normality summary for all rounds', str_list,font_size)    
    
        # create a notation page
    out_slide_layout = prs.slide_layouts[4]
    # Create the summary graph
    out_slide = prs.slides.add_slide(out_slide_layout)
    title = out_slide.shapes.title
    title.text = "Notations"
    subtitle = out_slide.placeholders[11]
    subtitle.text = "PBJ Rados Benchmark Result Summary"  
    subtitle = out_slide.placeholders[10]    
    tf = subtitle.text_frame 
    tf.text=''    
    font_size=16  
    
    str_list=['One to one comparison with 4 notations; the order of rounds is listed before',
              '1: two samples likely have the same mean',
              '0: two sample unlikely have no same mean',
              '-1: at least one vector has no enough data',
              '-2: at least one vector is not normal distribution']   
    add_items_to_tf(out_slide, 'Summary for pairs (one-way ANOVA for normal distribution)', str_list,font_size)   
    
    str_list=['One to one comparison with 3 notations; the order of rounds is listed before',
              '1: two samples likely from the same distribution',
              '0: two samples unlikely from the same distribution',
              '-1: at least one sample has no enough data']   
    add_items_to_tf(out_slide, 'Summary for pairs (two-sample KS test for continuous distribution)', str_list,font_size)   
    
    str_list=['One to one comparison with 3 notations; the order of rounds is listed before',
              '1: two samples likely have the same mean',
              '0: two samples unlikely have the same mean',
              '-1: at least one sample has no enough data']   
    add_items_to_tf(out_slide, 'Summary for pairs (two-sample T-test for identical mean)', str_list,font_size)  
    
    ## add the individual slides
    metric_str=['radosBench','radosBench_BW','Throughput (MBPS)']
    raw_data=add_metric_to_slide(prs,round_list,path0,metric_str,options,gen_sum=0,reduced_num=0, method_v=0)      
    
    metric_str=['radosBench','radosBench_LAT','Latency (second)']
    raw_data=add_metric_to_slide(prs,round_list,path0,metric_str,options,gen_sum=0,reduced_num=0, method_v=1)    
    
    prs.save('rados_compare2.pptx')
    
    return raw_data
    
########### main test####################

#round_list=['pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-10_170401033958.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-11_170401035608.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-12_170401041212.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-13_170401042817.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-14_170401044420.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-15_170401050021.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-16_170401051627.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-17_170401053227.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-18_170401054832.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-19_170401060433.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-1_170401011547.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-20_170401062038.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-21_170403095513.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-22_170403101127.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-23_170403102730.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-24_170403104330.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-25_170403105933.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-26_170403111530.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-27_170403113127.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-28_170403114732.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-29_170403120339.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-2_170401013150.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-30_170403121938.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-3_170401014747.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-4_170401020352.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-5_170401021952.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-6_170401023557.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-7_170401025200.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-8_170401030801.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-9_170401032402.json' ]    

#round_list=['pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-10_170401033958.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-11_170401035608.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-12_170401041212.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-14_170401044420.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-15_170401050021.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-16_170401051627.json',
#    'pbj_IWasp_BS_XeonE_11_OSDs_30G,10G_+_MVA72_1_OSD_10G,10G_2x.6Ghz_2G_MR-13_170401042817.json']
#
#path0='./json/'
#options=[0,0,0,0,0,1]    
#
#generate_report(round_list,'./json/',[1,1,1,1,1,1])    
#generate_report([],'./json2/',[1,1,1,1,1,1])    
#generate_report([],'./json2/')